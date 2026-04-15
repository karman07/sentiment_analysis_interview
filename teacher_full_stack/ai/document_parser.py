"""
document_parser.py — Extract plain text from PDF, DOCX, PPTX, XLSX, TXT, MD, CSV
"""
import os
import csv
import io
from pathlib import Path
from typing import Optional


def parse_document(file_path: str) -> str:
    """
    Parse a document at `file_path` and return its full text content.
    Raises ValueError if the file type is unsupported.
    """
    path = Path(file_path)
    if not path.exists():
        raise FileNotFoundError(f"File not found: {file_path}")

    ext = path.suffix.lower()

    if ext == ".pdf":
        return _parse_pdf(file_path)
    elif ext in (".docx",):
        return _parse_docx(file_path)
    elif ext in (".pptx",):
        return _parse_pptx(file_path)
    elif ext in (".xlsx",):
        return _parse_xlsx(file_path)
    elif ext in (".txt", ".md"):
        return path.read_text(encoding="utf-8", errors="replace")
    elif ext == ".csv":
        return _parse_csv(file_path)
    else:
        raise ValueError(f"Unsupported file type: {ext}")


def _parse_pdf(file_path: str) -> str:
    from pypdf import PdfReader
    reader = PdfReader(file_path)
    texts = []
    for page in reader.pages:
        text = page.extract_text()
        if text:
            texts.append(text.strip())
    return "\n\n".join(texts)


def _parse_docx(file_path: str) -> str:
    from docx import Document
    doc = Document(file_path)
    paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
    # Also extract tables
    for table in doc.tables:
        for row in table.rows:
            row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
            if row_text:
                paragraphs.append(row_text)
    return "\n\n".join(paragraphs)


def _parse_pptx(file_path: str) -> str:
    from pptx import Presentation
    prs = Presentation(file_path)
    slides_text = []
    for slide_num, slide in enumerate(prs.slides, 1):
        texts = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
        if texts:
            slides_text.append(f"[Slide {slide_num}]\n" + "\n".join(texts))
    return "\n\n".join(slides_text)


def _parse_xlsx(file_path: str) -> str:
    from openpyxl import load_workbook
    wb = load_workbook(file_path, read_only=True, data_only=True)
    sheets_text = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows = []
        for row in ws.iter_rows(values_only=True):
            cells = [str(cell) for cell in row if cell is not None]
            if cells:
                rows.append(" | ".join(cells))
        if rows:
            sheets_text.append(f"[Sheet: {sheet_name}]\n" + "\n".join(rows))
    wb.close()
    return "\n\n".join(sheets_text)


def _parse_csv(file_path: str) -> str:
    rows = []
    with open(file_path, encoding="utf-8", errors="replace") as f:
        reader = csv.reader(f)
        for row in reader:
            rows.append(" | ".join(row))
    return "\n".join(rows)
